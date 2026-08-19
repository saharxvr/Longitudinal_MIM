"""Build a PowerPoint with prior/current + 5-doctor collages, ordered by disagreement level.

For every pair, the deck contains three slides:
  1. Case header (e.g. "Case 12 - Disagreement level 5")
  2. Prior/Current image from per_pair_collages/pair{N}_prior_current.png
  3. 5 physician annotation collage from
     per_pair_collages_doctors_by_disagreement/disagreement_level_K/pair{N}_doctors.png

Pairs are ordered first by disagreement level (low to high by default, so
strongest agreement first), then by pair number.
"""
from __future__ import annotations

import argparse
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


_LEVEL_DIR_RE = re.compile(r"^disagreement_level_(\d+)$", re.IGNORECASE)
_DOCTOR_FILE_RE = re.compile(r"^pair(\d+)_doctors\.png$", re.IGNORECASE)


def _fit_rect(img_w: int, img_h: int, box_w: int, box_h: int) -> tuple[int, int]:
    if img_w <= 0 or img_h <= 0:
        return box_w, box_h
    scale = min(box_w / float(img_w), box_h / float(img_h))
    return max(1, int(round(img_w * scale))), max(1, int(round(img_h * scale)))


def _add_image_slide(prs: Presentation, image_path: Path, margin_in: float) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_w = prs.slide_width
    slide_h = prs.slide_height
    margin = Inches(margin_in)

    avail_x = margin
    avail_y = margin
    avail_w = slide_w - 2 * margin
    avail_h = slide_h - 2 * margin

    with Image.open(image_path) as im:
        img_w, img_h = im.size

    fit_w, fit_h = _fit_rect(img_w, img_h, int(avail_w), int(avail_h))
    left = int(avail_x + (avail_w - fit_w) / 2)
    top = int(avail_y + (avail_h - fit_h) / 2)
    slide.shapes.add_picture(str(image_path), left, top, width=fit_w, height=fit_h)


def _add_case_header_slide(prs: Presentation, pair_num: int, level: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    box_w = int(slide_w * 0.9)
    box_h = int(slide_h * 0.35)
    left = int((slide_w - box_w) / 2)
    top = int((slide_h - box_h) / 2)

    tx = slide.shapes.add_textbox(left, top, box_w, box_h)
    tf = tx.text_frame
    tf.clear()
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.alignment = 1  # center
    run = p.add_run()
    run.text = f"Case {pair_num}"
    run.font.size = Pt(72)
    run.font.bold = True

    p2 = tf.add_paragraph()
    p2.alignment = 1
    run2 = p2.add_run()
    run2.text = f"Disagreement level {level}"
    run2.font.size = Pt(40)


def _collect_doctor_collages(doctors_root: Path) -> dict[int, tuple[int, Path]]:
    """Return {pair_num: (level, doctor_collage_path)} for the doctors-by-level folder."""
    mapping: dict[int, tuple[int, Path]] = {}
    if not doctors_root.exists():
        raise FileNotFoundError(f"Doctors collages root not found: {doctors_root}")
    for level_dir in sorted(doctors_root.iterdir()):
        if not level_dir.is_dir():
            continue
        m = _LEVEL_DIR_RE.match(level_dir.name)
        if not m:
            continue
        level = int(m.group(1))
        for f in sorted(level_dir.iterdir()):
            mf = _DOCTOR_FILE_RE.match(f.name)
            if not mf:
                continue
            pair_num = int(mf.group(1))
            mapping[pair_num] = (level, f)
    return mapping


def _collect_prior_current(prior_current_root: Path) -> dict[int, Path]:
    """Return {pair_num: pairN_prior_current.png path}."""
    mapping: dict[int, Path] = {}
    if not prior_current_root.exists():
        raise FileNotFoundError(f"Prior/current collages root not found: {prior_current_root}")
    for f in sorted(prior_current_root.glob("pair*_prior_current.png")):
        m = re.match(r"^pair(\d+)_prior_current\.png$", f.name, re.IGNORECASE)
        if m:
            mapping[int(m.group(1))] = f
    return mapping


def build_presentation(args: argparse.Namespace) -> tuple[int, int]:
    prior_current_root = Path(args.prior_current_dir)
    doctors_root = Path(args.doctors_dir)
    out_path = Path(args.out)

    prior_current_by_pair = _collect_prior_current(prior_current_root)
    doctors_by_pair = _collect_doctor_collages(doctors_root)

    common_pairs = set(prior_current_by_pair) & set(doctors_by_pair)

    # Sort by (level, pair) so disagreement ordering is global.
    if args.descending:
        sorter = lambda pn: (-doctors_by_pair[pn][0], pn)
    else:
        sorter = lambda pn: (doctors_by_pair[pn][0], pn)
    ordered_pairs = sorted(common_pairs, key=sorter)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    pair_count = 0
    slide_count = 0

    for pair_num in ordered_pairs:
        level, doctor_path = doctors_by_pair[pair_num]
        prior_current_path = prior_current_by_pair[pair_num]

        _add_case_header_slide(prs, pair_num, level)
        _add_image_slide(prs, prior_current_path, args.margin_in)
        _add_image_slide(prs, doctor_path, args.margin_in)

        pair_count += 1
        slide_count += 3

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))

    return pair_count, slide_count


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description=__doc__)
    p.add_argument(
        "--prior-current-dir",
        type=Path,
        default=Path("Sahar_work/files/ov_results_sq/no_cc_itamar_segs_98/per_pair_collages"),
    )
    p.add_argument(
        "--doctors-dir",
        type=Path,
        default=Path(
            "Sahar_work/files/ov_results_sq/no_cc_itamar_segs_98/per_pair_collages_doctors_by_disagreement"
        ),
    )
    p.add_argument(
        "--out",
        type=Path,
        default=Path(
            "Sahar_work/files/ov_results_sq/no_cc_itamar_segs_98/pair_sequence_by_disagreement.pptx"
        ),
    )
    p.add_argument(
        "--descending",
        action="store_true",
        help="Order pairs from highest disagreement level to lowest (default: low to high)",
    )
    p.add_argument("--margin-in", type=float, default=0.1)
    return p.parse_args()


def main() -> int:
    args = parse_args()
    pairs, slides = build_presentation(args)
    print(f"Included pairs: {pairs}")
    print(f"Total slides:   {slides}")
    print(f"Wrote presentation: {args.out.resolve()}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
