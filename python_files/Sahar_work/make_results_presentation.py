from __future__ import annotations

import argparse
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from PIL import Image


TABLE_EQUATIONS = {
    "agreement_summary_human_human": (
        "Mean: μ = (1/N) Σ xᵢ\n"
        "SD: σ = sqrt((1/N) Σ (xᵢ-μ)²)\n"
        "IQR: Q3 - Q1\n"
        "Range: max(x) - min(x)\n"
        "95% CI (bootstrap mean): [Q2.5%(μ*), Q97.5%(μ*)]"
    ),
    "annotation_burden_by_physician": (
        "For physician r and type t:\n"
        "Total₍r,t₎ = Σ c₍r,t,ℓ₎\n"
        "Mean₍r,t₎ = (1/L) Σ c₍r,t,ℓ₎\n"
        "SD₍r,t₎ = sqrt((1/L) Σ (c₍r,t,ℓ₎-Mean₍r,t₎)²)\n"
        "Min/Max = min/max over ℓ"
    ),
    "consensus_sensitivity_summary": (
        "Sensitivity(k,t) = Detected(k,t) / Total(k,t)\n"
        "k: consensus level (1..5)\n"
        "t: positive or negative"
    ),
    "positive_vs_negative_tests": (
        "Δ = (1/N) Σ (xᵢ - yᵢ)\n"
        "Exact sign-flip p-value:\n"
        "p = (1/2ᴺ) Σ 1{|(1/N) Σ sᵢ(xᵢ-yᵢ)| ≥ |Δ|}\n"
        "sᵢ ∈ {-1,+1}"
    ),
    "raw_hh_pair_values": (
        "Raw inputs used by summaries/tests:\n"
        "{xᵢ} for each condition\n"
        "i = 1..10 physician pairs"
    ),
}


def fit_rect(img_w: int, img_h: int, box_w: int, box_h: int) -> tuple[int, int]:
    scale = min(box_w / float(img_w), box_h / float(img_h))
    return max(1, int(round(img_w * scale))), max(1, int(round(img_h * scale)))


def add_image_slide(prs: Presentation, image_path: Path, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    margin = Inches(0.25)
    title_h = Inches(0.55)

    tx = slide.shapes.add_textbox(margin, Inches(0.05), slide_w - 2 * margin, title_h)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(30)
    run.font.bold = True

    avail_x = margin
    avail_y = title_h + margin
    avail_w = slide_w - 2 * margin
    avail_h = slide_h - title_h - 2 * margin

    with Image.open(image_path) as im:
        iw, ih = im.size
    fw, fh = fit_rect(iw, ih, int(avail_w), int(avail_h))
    left = int(avail_x + (avail_w - fw) / 2)
    top = int(avail_y + (avail_h - fh) / 2)

    slide.shapes.add_picture(str(image_path), left, top, width=fw, height=fh)


def add_table_with_equation_slide(prs: Presentation, image_path: Path, title: str, equation_text: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    margin = Inches(0.2)
    title_h = Inches(0.55)
    gap = Inches(0.2)

    tx = slide.shapes.add_textbox(margin, Inches(0.05), slide_w - 2 * margin, title_h)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True

    content_y = title_h + margin
    content_h = slide_h - title_h - 2 * margin

    left_w = int((slide_w - 2 * margin - gap) * 0.72)
    right_w = int((slide_w - 2 * margin - gap) * 0.28)

    left_x = margin
    right_x = left_x + left_w + gap

    with Image.open(image_path) as im:
        iw, ih = im.size
    fw, fh = fit_rect(iw, ih, left_w, int(content_h))
    img_left = int(left_x + (left_w - fw) / 2)
    img_top = int(content_y + (content_h - fh) / 2)
    slide.shapes.add_picture(str(image_path), img_left, img_top, width=fw, height=fh)

    eq_box = slide.shapes.add_textbox(right_x, content_y, right_w, content_h)
    etf = eq_box.text_frame
    etf.clear()
    etf.word_wrap = True
    lines = equation_text.split("\n")
    if lines:
        p0 = etf.paragraphs[0]
        p0.text = lines[0]
        p0.font.size = Pt(18)
        p0.font.bold = True
        for line in lines[1:]:
            p2 = etf.add_paragraph()
            p2.text = line
            p2.font.size = Pt(16)


def add_editable_table_with_equation_slide(prs: Presentation, csv_path: Path, title: str, equation_text: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    margin = Inches(0.2)
    title_h = Inches(0.55)
    gap = Inches(0.2)

    tx = slide.shapes.add_textbox(margin, Inches(0.05), slide_w - 2 * margin, title_h)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(28)
    run.font.bold = True

    content_y = title_h + margin
    content_h = slide_h - title_h - 2 * margin

    left_w = int((slide_w - 2 * margin - gap) * 0.72)
    right_w = int((slide_w - 2 * margin - gap) * 0.28)

    left_x = margin
    right_x = left_x + left_w + gap

    df = pd.read_csv(csv_path)
    nrows, ncols = df.shape
    total_rows = nrows + 1

    table_shape = slide.shapes.add_table(total_rows, ncols, left_x, content_y, left_w, content_h)
    table = table_shape.table

    # Dynamic column widths by content length.
    max_lens = []
    for c, col_name in enumerate(df.columns):
        col_vals = [str(v) for v in df.iloc[:, c].tolist()]
        m = max([len(str(col_name))] + [len(v) for v in col_vals])
        max_lens.append(max(4, m))
    total_len = sum(max_lens)
    for c in range(ncols):
        ratio = max_lens[c] / total_len
        table.columns[c].width = int(left_w * ratio)

    # Header row.
    for c, col_name in enumerate(df.columns):
        cell = table.cell(0, c)
        cell.text = str(col_name)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(58, 95, 138)
        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(255, 255, 255)

    # Data rows.
    for r in range(nrows):
        for c in range(ncols):
            cell = table.cell(r + 1, c)
            val = df.iat[r, c]
            if isinstance(val, float):
                txt = f"{val:.4f}"
            else:
                txt = str(val)
            cell.text = txt
            cell.fill.solid()
            if (r + 1) % 2 == 0:
                cell.fill.fore_color.rgb = RGBColor(243, 246, 250)
            else:
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)

    eq_box = slide.shapes.add_textbox(right_x, content_y, right_w, content_h)
    etf = eq_box.text_frame
    etf.clear()
    etf.word_wrap = True
    lines = equation_text.split("\n")
    if lines:
        p0 = etf.paragraphs[0]
        p0.text = lines[0]
        p0.font.size = Pt(18)
        p0.font.bold = True
        for line in lines[1:]:
            p2 = etf.add_paragraph()
            p2.text = line
            p2.font.size = Pt(16)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    box_w = int(slide_w * 0.9)
    left = int((slide_w - box_w) / 2)

    tx1 = slide.shapes.add_textbox(left, int(slide_h * 0.30), box_w, int(slide_h * 0.2))
    t1 = tx1.text_frame
    t1.clear()
    p1 = t1.paragraphs[0]
    p1.alignment = 1
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(44)
    r1.font.bold = True

    tx2 = slide.shapes.add_textbox(left, int(slide_h * 0.50), box_w, int(slide_h * 0.15))
    t2 = tx2.text_frame
    t2.clear()
    p2 = t2.paragraphs[0]
    p2.alignment = 1
    r2 = p2.add_run()
    r2.text = subtitle
    r2.font.size = Pt(24)


def main() -> int:
    parser = argparse.ArgumentParser(description="Create PPT presentation from paper_stats tables and plots.")
    parser.add_argument(
        "--paper-stats-dir",
        type=Path,
        default=Path("Sahar_work/files/ov_results_sq/no_cc_itamar_segs_98/paper_stats"),
    )
    parser.add_argument(
        "--out",
        type=Path,
        default=Path("Sahar_work/files/ov_results_sq/no_cc_itamar_segs_98/paper_stats/observer_variability_results_presentation.pptx"),
    )
    args = parser.parse_args()

    stats_dir = args.paper_stats_dir
    tables_dir = stats_dir / "tables"
    plots_dir = stats_dir / "plots"

    csv_files = [
        tables_dir / "agreement_summary_human_human.csv",
        tables_dir / "annotation_burden_by_physician.csv",
        tables_dir / "consensus_sensitivity_summary.csv",
        tables_dir / "positive_vs_negative_tests.csv",
        tables_dir / "raw_hh_pair_values.csv",
    ]

    plot_files = [
        plots_dir / "heatmap_per_pair_positive_hh.png",
        plots_dir / "heatmap_per_pair_negative_hh.png",
        plots_dir / "heatmap_per_pair_all_hh.png",
        plots_dir / "heatmap_per_detection_positive_hh.png",
        plots_dir / "heatmap_per_detection_negative_hh.png",
        plots_dir / "heatmap_per_detection_all_hh.png",
        plots_dir / "boxplot_pairwise_hh.png",
        plots_dir / "boxplot_per_detection_hh.png",
        plots_dir / "consensus_sensitivity_positive_negative.png",
        plots_dir / "annotation_burden_by_physician.png",
    ]

    for csv_path in csv_files:
        if not csv_path.exists():
            raise FileNotFoundError(f"Missing table file: {csv_path}")
    for plot_path in plot_files:
        if not plot_path.exists():
            raise FileNotFoundError(f"Missing plot file: {plot_path}")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        "Observer Variability Results",
        "98 pairs, 5 physicians | Human-human analysis",
    )

    add_title_slide(prs, "Summary Tables", "Statistical summaries and tests")
    for csv_path in csv_files:
        add_editable_table_with_equation_slide(
            prs,
            csv_path,
            csv_path.stem.replace("_", " ").title(),
            TABLE_EQUATIONS.get(csv_path.stem, "Equation not available"),
        )

    add_title_slide(prs, "Agreement Heatmaps", "Pairwise and per-detection agreement")
    for p in plot_files[:6]:
        add_image_slide(prs, p, p.stem.replace("_", " ").title())

    add_title_slide(prs, "Distribution and Trend Plots", "Agreement distributions and consensus trend")
    for p in plot_files[6:]:
        add_image_slide(prs, p, p.stem.replace("_", " ").title())

    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.out))

    print(f"Presentation created: {args.out.resolve()}")
    print(f"Slides: {len(prs.slides)}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
