from __future__ import annotations

import argparse
import json
from dataclasses import dataclass
from pathlib import Path

import matplotlib.pyplot as plt
from matplotlib import colors
from matplotlib.patches import Ellipse
import nibabel as nib
import numpy as np
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


@dataclass(frozen=True)
class EllipseAnn:
    cx: float
    cy: float
    rx: float
    ry: float
    angle_deg: float
    label: str
    tag: str


def load_ellipses(json_path: Path) -> list[EllipseAnn]:
    if not json_path.exists():
        return []
    data = json.loads(json_path.read_text(encoding="utf-8-sig"))
    if not isinstance(data, list):
        return []

    out: list[EllipseAnn] = []
    for item in data[1:]:
        if not isinstance(item, dict):
            continue
        out.append(
            EllipseAnn(
                cx=float(item.get("cx", 0.0)),
                cy=float(item.get("cy", 0.0)),
                rx=float(item.get("rx", 0.0)),
                ry=float(item.get("ry", 0.0)),
                angle_deg=float(item.get("angle", 0.0)),
                label=str(item.get("label", "")),
                tag=str(item.get("tag", "")),
            )
        )
    return out


def load_current_image(pair_dir: Path, out_size: int = 792) -> np.ndarray:
    nii_files = sorted([p for p in pair_dir.glob("*.nii.gz") if not p.name.endswith("_lung_seg.nii.gz")])
    if len(nii_files) < 2:
        raise FileNotFoundError(f"Expected 2 scans in {pair_dir}, found {len(nii_files)}")

    current = nib.load(str(nii_files[1])).get_fdata().T.astype(np.float32)
    if current.ndim == 3:
        current = current[:, :, current.shape[2] // 2]

    mn = float(np.min(current))
    mx = float(np.max(current))
    if mx > mn:
        current = (current - mn) / (mx - mn)
    else:
        current = np.zeros_like(current, dtype=np.float32)

    if current.shape != (out_size, out_size):
        img = Image.fromarray((current * 255.0).astype(np.uint8), mode="L")
        current = np.asarray(img.resize((out_size, out_size), Image.Resampling.BILINEAR), dtype=np.float32) / 255.0

    return current


def load_model_map(pred_pair_dir: Path, out_size: int = 792) -> np.ndarray:
    out_path = pred_pair_dir / "output.nii.gz"
    if not out_path.exists():
        return np.zeros((out_size, out_size), dtype=np.float32)

    arr = nib.load(str(out_path)).get_fdata().T.astype(np.float32)
    if arr.ndim == 3:
        arr = arr[:, :, arr.shape[2] // 2]

    if arr.shape != (out_size, out_size):
        arr_img = Image.fromarray(arr, mode="F")
        arr = np.asarray(arr_img.resize((out_size, out_size), Image.Resampling.BILINEAR), dtype=np.float32)

    return arr


def pred_cmap() -> colors.LinearSegmentedColormap:
    return colors.LinearSegmentedColormap.from_list(
        "prediction_diff",
        (
            (0.000, (0.235, 1.000, 0.239)),
            (0.400, (0.000, 1.000, 0.702)),
            (0.500, (1.000, 0.988, 0.988)),
            (0.600, (1.000, 0.604, 0.000)),
            (1.000, (0.682, 0.000, 0.000)),
        ),
    )


def alpha_map(x: np.ndarray) -> np.ndarray:
    x_abs = np.abs(x)
    max_val = max(float(np.max(x_abs)), 0.05)
    return x_abs / max_val


def color_for_label(label: str) -> tuple[float, float, float]:
    m = {
        "Appearance": (1.0, 0.0, 1.0),
        "Disappearance": (0.2, 0.6, 1.0),
        "Persistence": (1.0, 1.0, 0.0),
    }
    return m.get(label, (1.0, 1.0, 1.0))


def draw_ellipses(ax: plt.Axes, ellipses: list[EllipseAnn], annotate: bool = False) -> None:
    for e in ellipses:
        patch = Ellipse(
            (e.cx, e.cy),
            width=2.0 * e.rx,
            height=2.0 * e.ry,
            angle=e.angle_deg,
            fill=False,
            linewidth=2.0,
            edgecolor=color_for_label(e.label),
        )
        ax.add_patch(patch)
        if annotate:
            ax.text(
                e.cx - e.rx,
                e.cy - e.ry - 8,
                f"{e.label} | {e.tag}" if e.tag else e.label,
                fontsize=6,
                color="white",
                bbox={"facecolor": "black", "alpha": 0.5, "pad": 1},
            )


def make_pair_figure(
    pair: str,
    current_img: np.ndarray,
    model_map: np.ndarray,
    n_ellipses: list[EllipseAnn],
    t_ellipses: list[EllipseAnn],
    out_path: Path,
) -> None:
    fig, axs = plt.subplots(2, 2, figsize=(12, 10), dpi=150)
    cm = pred_cmap()
    divnorm = colors.TwoSlopeNorm(
        vmin=min(float(np.min(model_map)), -0.01),
        vcenter=0.0,
        vmax=max(float(np.max(model_map)), 0.01),
    )

    ax = axs[0, 0]
    ax.imshow(current_img, cmap="gray", origin="upper")
    ax.imshow(model_map, cmap=cm, norm=divnorm, alpha=alpha_map(model_map), origin="upper")
    draw_ellipses(ax, n_ellipses, annotate=False)
    ax.set_title(f"{pair} - Nathalia + Model")
    ax.axis("off")

    ax = axs[0, 1]
    ax.imshow(current_img, cmap="gray", origin="upper")
    ax.imshow(model_map, cmap=cm, norm=divnorm, alpha=alpha_map(model_map), origin="upper")
    draw_ellipses(ax, t_ellipses, annotate=False)
    ax.set_title(f"{pair} - Thaer + Model")
    ax.axis("off")

    ax = axs[1, 0]
    ax.imshow(current_img, cmap="gray", origin="upper")
    draw_ellipses(ax, n_ellipses, annotate=False)
    draw_ellipses(ax, t_ellipses, annotate=False)
    ax.set_title(f"{pair} - Nathalia + Thaer")
    ax.axis("off")

    ax = axs[1, 1]
    ax.imshow(current_img, cmap="gray", origin="upper")
    im = ax.imshow(model_map, cmap=cm, norm=divnorm, alpha=alpha_map(model_map), origin="upper")
    ax.set_title(f"{pair} - Model Output")
    ax.axis("off")
    fig.colorbar(im, ax=ax, fraction=0.046, pad=0.04)

    fig.tight_layout()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path)
    plt.close(fig)


def make_rates_heatmap(summary: dict, out_path: Path) -> None:
    rates = summary["rates"]
    mat = np.array(
        [
            [1.0, rates["n_vs_t"], rates["n_vs_m"]],
            [rates["t_vs_n"], 1.0, rates["t_vs_m"]],
            [rates["m_vs_n"], rates["m_vs_t"], 1.0],
        ],
        dtype=np.float64,
    )
    labels = ["Nathalia", "Thaer", "Model"]

    fig, ax = plt.subplots(figsize=(6, 5), dpi=160)
    im = ax.imshow(mat, cmap="viridis", vmin=0.0, vmax=1.0)
    ax.set_xticks(np.arange(3), labels)
    ax.set_yticks(np.arange(3), labels)

    for i in range(3):
        for j in range(3):
            ax.text(j, i, f"{mat[i, j]:.2f}", ha="center", va="center", color="white", fontsize=11)

    ax.set_title("Mini OV Pairwise Agreement Rates")
    fig.colorbar(im, fraction=0.046, pad=0.04)
    fig.tight_layout()
    out_path.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(out_path)
    plt.close(fig)


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_summary_slide(
    prs: Presentation,
    run_summary: dict,
    sensitivity: dict,
    precision: dict,
    onepx_summary: dict,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "OV Summary (ICU-Style)"

    lines = [
        f"Pairs processed: {run_summary['pairs_processed']}",
        f"Observers: {', '.join(run_summary['observers'])}",
        f"Sensitivity L1 (Pos/Neg): {sensitivity['Sensitivity Consensus Level 1 (Positive)']:.3f} / {sensitivity['Sensitivity Consensus Level 1 (Negative)']:.3f}",
        f"Sensitivity L2 (Pos/Neg): {sensitivity['Sensitivity Consensus Level 2 (Positive)']:.3f} / {sensitivity['Sensitivity Consensus Level 2 (Negative)']:.3f}",
        f"HMDR Nathalia (Pos/Neg): {precision['Nathalia HMDR (Positive)']:.3f} / {precision['Nathalia HMDR (Negative)']:.3f}",
        f"HMDR Thaer (Pos/Neg): {precision['Thaer HMDR (Positive)']:.3f} / {precision['Thaer HMDR (Negative)']:.3f}",
        f"HMDR Model (Pos/Neg): {precision['Model HMDR (Positive)']:.3f} / {precision['Model HMDR (Negative)']:.3f}",
        f"UDPP Model (Pos/Neg): {precision['UDPP Model (Positive)']:.3f} / {precision['UDPP Model (Negative)']:.3f}",
        f"Legacy doc-doc combined agree/disagree: {onepx_summary['combined_directional']['agree']} / {onepx_summary['combined_directional']['disagree']}",
    ]

    box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(12.0), Inches(5.6))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0


def add_figure_slide(prs: Presentation, title: str, fig_path: Path) -> None:
    if not fig_path.exists():
        return
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    slide.shapes.add_picture(str(fig_path), Inches(0.4), Inches(1.0), width=Inches(12.5))


def add_pairs_table_slide(prs: Presentation, title: str, rows: list[dict]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title

    headers = ["Pair", "N dis", "T dis", "M comps", "M pixels"]
    table = slide.shapes.add_table(rows=len(rows) + 1, cols=len(headers), left=Inches(0.6), top=Inches(1.3), width=Inches(12.0), height=Inches(5.6)).table

    for c, h in enumerate(headers):
        table.cell(0, c).text = h

    for r, row in enumerate(rows, start=1):
        table.cell(r, 0).text = str(row.get("pair", ""))
        table.cell(r, 1).text = str(row.get("n_disagree", row.get("n_vs_t_disagree", 0)))
        table.cell(r, 2).text = str(row.get("t_disagree", row.get("t_vs_n_disagree", 0)))
        table.cell(r, 3).text = str(row.get("m_components_total", 0))
        table.cell(r, 4).text = str(row.get("model_pixels", 0))


def add_pair_figure_slide(prs: Presentation, pair: str, fig_path: Path) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = f"Pair Detail: {pair}"
    slide.shapes.add_picture(str(fig_path), Inches(0.4), Inches(1.0), width=Inches(12.5))


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Create presentation from ICU-style OV outputs and pair visuals for Nathalia, Thaer, and model output."
    )
    parser.add_argument(
        "--json-root",
        type=Path,
        default=Path("python_files") / "annotation tool" / "Annotations_Pnimit" / "JSON_only_named",
    )
    parser.add_argument(
        "--pairs-root",
        type=Path,
        default=Path("python_files") / "annotation tool" / "Pairs_PNIMIT_1_pairs",
    )
    parser.add_argument(
        "--predictions-root",
        type=Path,
        default=Path("python_files") / "annotation tool" / "predictions_pnimit_lung",
    )
    parser.add_argument(
        "--ov-dir",
        type=Path,
        default=Path("python_files")
        / "annotation tool"
        / "Annotations_Pnimit"
        / "JSON_only_named"
        / "ov_sq_style_nathalia_thaer_model",
        help="Directory containing ICU-style OV outputs.",
    )
    parser.add_argument(
        "--onepx-summary",
        type=Path,
        default=Path("python_files")
        / "annotation tool"
        / "Annotations_Pnimit"
        / "JSON_only_named"
        / "thaer_nathalia_agreement_1px_summary.json",
    )
    parser.add_argument(
        "--out-dir",
        type=Path,
        default=Path("python_files")
        / "annotation tool"
        / "Annotations_Pnimit"
        / "JSON_only_named"
        / "presentation_nathalia_thaer_model",
    )
    parser.add_argument(
        "--top-k",
        type=int,
        default=8,
        help="Number of highest-disagreement pairs to include as detailed slides.",
    )
    return parser.parse_args()


def run(args: argparse.Namespace) -> int:
    out_dir = args.out_dir
    figs_dir = out_dir / "figures"
    out_dir.mkdir(parents=True, exist_ok=True)
    figs_dir.mkdir(parents=True, exist_ok=True)

    run_summary = json.loads((args.ov_dir / "run_summary.json").read_text(encoding="utf-8"))
    sensitivity = json.loads((args.ov_dir / "sensitivity_measures.json").read_text(encoding="utf-8"))
    precision = json.loads((args.ov_dir / "precision_measures.json").read_text(encoding="utf-8"))
    onepx_data = json.loads(args.onepx_summary.read_text(encoding="utf-8"))

    per_pair = list(onepx_data["per_pair"])

    per_pair_sorted = sorted(
        per_pair,
        key=lambda r: int(r.get("n_disagree", 0)) + int(r.get("t_disagree", 0)),
        reverse=True,
    )
    top_rows = per_pair_sorted[: max(int(args.top_k), 1)]

    n_dir = args.json_root / "Nathalia"
    t_dir = args.json_root / "Thaer"

    figure_paths: list[tuple[str, Path]] = []
    for row in top_rows:
        pair_raw = str(row["pair"])
        pair = pair_raw[:-5] if pair_raw.endswith(".json") else pair_raw
        n_json = n_dir / f"Nathalia_{pair}.json"
        t_json = t_dir / f"Thaer_{pair}.json"
        pair_dir = args.pairs_root / pair
        pred_dir = args.predictions_root / pair

        if not pair_dir.exists():
            continue

        try:
            current = load_current_image(pair_dir, out_size=792)
            model = load_model_map(pred_dir, out_size=792)
            n_ells = load_ellipses(n_json)
            t_ells = load_ellipses(t_json)
        except Exception:
            continue

        fig_path = figs_dir / f"{pair}_overview.png"
        make_pair_figure(pair, current, model, n_ells, t_ells, fig_path)
        figure_paths.append((pair, fig_path))

    prs = Presentation()
    add_title_slide(
        prs,
        title="OV Report: Nathalia vs Thaer vs Model",
        subtitle="ICU-style OV metrics on PNIMIT pairs",
    )
    add_summary_slide(
        prs,
        run_summary=run_summary,
        sensitivity=sensitivity,
        precision=precision,
        onepx_summary=onepx_data["summary"],
    )
    add_figure_slide(prs, "PAI Per Label (All)", args.ov_dir / "per_label_agreement_all.png")
    add_figure_slide(prs, "PAI Per Pair (All)", args.ov_dir / "per_pair_agreement_all.png")
    add_figure_slide(prs, "Sensitivity Model", args.ov_dir / "sensitivity_consensus_levels_model.png")
    add_figure_slide(prs, "Sensitivity Nathalia", args.ov_dir / "sensitivity_consensus_levels_nathalia.png")
    add_figure_slide(prs, "Sensitivity Thaer", args.ov_dir / "sensitivity_consensus_levels_thaer.png")
    add_pairs_table_slide(prs, title="Top Pair Disagreements (Doctor-to-Doctor)", rows=top_rows[:10])

    for pair, fig_path in figure_paths:
        add_pair_figure_slide(prs, pair, fig_path)

    pptx_path = out_dir / "nathalia_thaer_model_ov_sq_style.pptx"
    prs.save(str(pptx_path))

    txt_summary_path = out_dir / "presentation_summary.txt"
    txt_summary_path.write_text(
        "\n".join(
            [
                "ICU-style OV presentation generated.",
                f"Slides: {len(prs.slides)}",
                f"Top pairs rendered: {len(figure_paths)}",
                f"PPTX: {pptx_path}",
                f"OV source: {args.ov_dir}",
            ]
        ),
        encoding="utf-8",
    )

    print(f"Saved: {pptx_path}")
    print(f"Saved: {txt_summary_path}")
    return 0


if __name__ == "__main__":
    raise SystemExit(run(parse_args()))
