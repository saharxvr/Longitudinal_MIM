from __future__ import annotations

import argparse
import re
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt


PAIR_RE = re.compile(r"^pair(\d+)_(prior_current|annotators_model)\.png$", re.IGNORECASE)


def _collect_pairs(collages_dir: Path) -> dict[int, dict[str, Path]]:
    pairs: dict[int, dict[str, Path]] = {}
    for p in collages_dir.glob("pair*_*.png"):
        m = PAIR_RE.match(p.name)
        if not m:
            continue
        pair_num = int(m.group(1))
        kind = m.group(2).lower()
        pairs.setdefault(pair_num, {})[kind] = p
    return pairs


def _fit_rect(img_w: int, img_h: int, box_w: int, box_h: int) -> tuple[int, int]:
    if img_w <= 0 or img_h <= 0:
        return box_w, box_h
    scale = min(box_w / float(img_w), box_h / float(img_h))
    return max(1, int(round(img_w * scale))), max(1, int(round(img_h * scale)))


def _add_image_slide(prs: Presentation, image_path: Path, title: str, margin_in: float, title_in: float) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_w = prs.slide_width
    slide_h = prs.slide_height

    margin = Inches(margin_in)

    avail_x = margin
    avail_y = margin
    avail_w = slide_w - 2 * margin
    avail_h = slide_h - 2 * margin

    with Image.open(image_path) as im:
        img_w, img_h = im.size

    fit_w, fit_h = _fit_rect(img_w, img_h, int(avail_w), int(avail_h))
    left = int(avail_x + (avail_w - fit_w) / 2)
    top = int(avail_y + (avail_h - fit_h) / 2)

    slide.shapes.add_picture(str(image_path), left, top, width=fit_w, height=fit_h)


def _add_case_header_slide(prs: Presentation, pair_num: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_w = prs.slide_width
    slide_h = prs.slide_height

    box_w = int(slide_w * 0.9)
    box_h = int(slide_h * 0.25)
    left = int((slide_w - box_w) / 2)
    top = int((slide_h - box_h) / 2)

    tx = slide.shapes.add_textbox(left, top, box_w, box_h)
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = 1  # center
    run = p.add_run()
    run.text = f"Case {pair_num}"
    run.font.size = Pt(72)
    run.font.bold = True


def build_presentation(collages_dir: Path, out_path: Path, margin_in: float, title_in: float) -> tuple[int, int]:
    prs = Presentation()
    # Force 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    pairs = _collect_pairs(collages_dir)
    ordered = sorted(pairs.keys())

    pair_count = 0
    slide_count = 0

    for pair_num in ordered:
        assets = pairs[pair_num]
        prior_current = assets.get("prior_current")
        ann_model = assets.get("annotators_model")
        if prior_current is None or ann_model is None:
            continue

        _add_case_header_slide(prs, pair_num)

        _add_image_slide(
            prs,
            prior_current,
            f"Pair {pair_num} - Prior and Current",
            margin_in=margin_in,
            title_in=title_in,
        )
        _add_image_slide(
            prs,
            ann_model,
            f"Pair {pair_num} - Physician Annotations and Model Output",
            margin_in=margin_in,
            title_in=title_in,
        )

        pair_count += 1
        slide_count += 3

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return pair_count, slide_count


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Create a sequential pair presentation from per-pair collages.")
    p.add_argument(
        "--collages-dir",
        type=Path,
        default=Path("Sahar_work/files/ov_results_sq/no_cc_itamar_segs_98/per_pair_collages"),
        help="Folder containing pairN_prior_current.png and pairN_annotators_model.png",
    )
    p.add_argument(
        "--out",
        type=Path,
        default=Path("Sahar_work/files/ov_results_sq/no_cc_itamar_segs_98/pair_sequence_presentation.pptx"),
        help="Output presentation path",
    )
    p.add_argument("--margin-in", type=float, default=0.1, help="Slide margin in inches")
    p.add_argument("--title-in", type=float, default=0.5, help="Title area height in inches")
    return p.parse_args()


def main() -> int:
    args = parse_args()
    if not args.collages_dir.exists():
        raise FileNotFoundError(f"Collages dir not found: {args.collages_dir}")

    pairs, slides = build_presentation(args.collages_dir, args.out, args.margin_in, args.title_in)
    print(f"Created presentation: {args.out.resolve()}")
    print(f"Included pairs: {pairs}")
    print(f"Total slides: {slides}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
